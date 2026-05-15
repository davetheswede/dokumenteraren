from __future__ import annotations

import html
import io
import json
import os
import posixpath
import re
import subprocess
import sys
import tempfile
import zipfile
from email.message import EmailMessage
from pathlib import Path


PROJECT_ROOT = Path(__file__).resolve().parents[1]
sys.path.insert(0, str(PROJECT_ROOT))


def configure_isolated_runtime() -> tempfile.TemporaryDirectory[str]:
    data_dir = tempfile.TemporaryDirectory(prefix="dokumenteraren-acceptance-")
    os.environ["DATA_DIR"] = data_dir.name
    os.environ["IMPORT_DIR"] = str(Path(data_dir.name).with_name(f"{Path(data_dir.name).name}-import"))
    os.environ["SESSION_SECRET"] = "acceptance-test-secret"
    os.environ["MAX_UPLOAD_BYTES"] = str(1_000_000)
    return data_dir


runtime_dir = configure_isolated_runtime()

from docx import Document as DocxDocument  # noqa: E402
from fastapi.testclient import TestClient  # noqa: E402
from odf.opendocument import OpenDocumentSpreadsheet, OpenDocumentText  # noqa: E402
from odf.table import Table, TableCell, TableRow  # noqa: E402
from odf.text import P  # noqa: E402
from openpyxl import Workbook  # noqa: E402
from PIL import Image  # noqa: E402
from pptx import Presentation  # noqa: E402

from app import db  # noqa: E402
from app.catalog import DOCUMENT_TEMPLATES  # noqa: E402
from app.crypto import ENC_PREFIX  # noqa: E402
from app.main import app  # noqa: E402
from app.security import redact_text  # noqa: E402
from app.services import importer, mail_importer  # noqa: E402
from app.services.export import safe_zip_member  # noqa: E402


CSRF_RE = re.compile(r'name="csrf_token" value="([^"]+)"')


def csrf(response) -> str:
    match = CSRF_RE.search(response.text)
    assert match, "CSRF-token saknas i formulär."
    return html.unescape(match.group(1))


def make_pdf(text: str) -> bytes:
    stream = f"BT /F1 18 Tf 72 720 Td ({text}) Tj ET".encode("latin-1")
    objects = [
        b"<< /Type /Catalog /Pages 2 0 R >>",
        b"<< /Type /Pages /Kids [3 0 R] /Count 1 >>",
        b"<< /Type /Page /Parent 2 0 R /MediaBox [0 0 612 792] /Resources << /Font << /F1 4 0 R >> >> /Contents 5 0 R >>",
        b"<< /Type /Font /Subtype /Type1 /BaseFont /Helvetica >>",
        b"<< /Length " + str(len(stream)).encode("ascii") + b" >>\nstream\n" + stream + b"\nendstream",
    ]
    output = io.BytesIO()
    output.write(b"%PDF-1.4\n")
    offsets = []
    for index, obj in enumerate(objects, start=1):
        offsets.append(output.tell())
        output.write(f"{index} 0 obj\n".encode("ascii"))
        output.write(obj)
        output.write(b"\nendobj\n")
    xref = output.tell()
    output.write(f"xref\n0 {len(objects) + 1}\n0000000000 65535 f \n".encode("ascii"))
    for offset in offsets:
        output.write(f"{offset:010d} 00000 n \n".encode("ascii"))
    output.write(
        f"trailer << /Root 1 0 R /Size {len(objects) + 1} >>\nstartxref\n{xref}\n%%EOF\n".encode("ascii")
    )
    return output.getvalue()


def make_docx(text: str) -> bytes:
    buffer = io.BytesIO()
    doc = DocxDocument()
    doc.add_paragraph(text)
    doc.save(buffer)
    return buffer.getvalue()


def make_xlsx(text: str) -> bytes:
    buffer = io.BytesIO()
    workbook = Workbook()
    sheet = workbook.active
    sheet.title = "Verifiering"
    sheet.append(["typ", "innehåll"])
    sheet.append(["xlsx", text])
    workbook.save(buffer)
    return buffer.getvalue()


def make_pptx(text: str) -> bytes:
    buffer = io.BytesIO()
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    text_box = slide.shapes.add_textbox(0, 0, 5_000_000, 1_000_000)
    text_box.text = text
    presentation.save(buffer)
    return buffer.getvalue()


def make_odt(text: str) -> bytes:
    buffer = io.BytesIO()
    doc = OpenDocumentText()
    doc.text.addElement(P(text=text))
    doc.save(buffer)
    return buffer.getvalue()


def make_ods(text: str) -> bytes:
    buffer = io.BytesIO()
    doc = OpenDocumentSpreadsheet()
    table = Table(name="Verifiering")
    row = TableRow()
    cell = TableCell()
    cell.addElement(P(text=text))
    row.addElement(cell)
    table.addElement(row)
    doc.spreadsheet.addElement(table)
    doc.save(buffer)
    return buffer.getvalue()


def make_flat_odf(text: str) -> bytes:
    return f"""<?xml version="1.0" encoding="UTF-8"?>
<office:document
  xmlns:office="urn:oasis:names:tc:opendocument:xmlns:office:1.0"
  xmlns:text="urn:oasis:names:tc:opendocument:xmlns:text:1.0"
  office:mimetype="application/vnd.oasis.opendocument.text">
  <office:body><office:text><text:p>{text}</text:p></office:text></office:body>
</office:document>
""".encode("utf-8")


def make_png() -> bytes:
    buffer = io.BytesIO()
    Image.new("RGB", (24, 16), color=(80, 120, 160)).save(buffer, format="PNG")
    return buffer.getvalue()


SAMPLES: list[tuple[str, bytes, str, bool]] = [
    ("text-token.txt", b"uniktext alfa person 19800101-1234\n", "uniktext", True),
    ("markdown-token.md", b"# Rubrik\nunikmarkdown beta\n", "unikmarkdown", True),
    ("forsakring.pdf", make_pdf("unikpdf krock villkor"), "unikpdf", True),
    ("word.docx", make_docx("unikdocx avtal"), "unikdocx", True),
    ("macro-word.docm", make_docx("unikdocm avtal"), "unikdocm", True),
    ("sheet.xlsx", make_xlsx("unikxlsx kvitto"), "unikxlsx", True),
    ("macro-sheet.xlsm", make_xlsx("unikxlsm kvitto"), "unikxlsm", True),
    ("presentation.pptx", make_pptx("unikpptx villkor"), "unikpptx", True),
    ("open.odt", make_odt("unikodt journal"), "unikodt", True),
    ("open.ods", make_ods("unikods kalkyl"), "unikods", True),
    ("flat.fodt", make_flat_odf("unikfodt anteckning"), "unikfodt", True),
    ("rich.rtf", b"{\\rtf1\\ansi unikrtf garanti}", "unikrtf", True),
    ("data.csv", b"kolumn\nunikcsv\n", "unikcsv", True),
    ("data.json", b'{"nyckel":"unikjson"}', "unikjson", True),
    ("data.yaml", b"nyckel: unikyaml\n", "unikyaml", True),
    ("mail.eml", b"Subject: Test\nFrom: a@example.test\nTo: b@example.test\n\nunikeml meddelande", "unikeml", True),
    ("photo.png", make_png(), "width", False),
]


def assert_zip_member_safe(name: str) -> None:
    assert not name.startswith(("/", "\\")), f"ZIP-namn är absolut: {name}"
    assert "\\" not in name, f"ZIP-namn innehåller backslash: {name}"
    assert ":" not in name.split("/", 1)[0], f"ZIP-namn innehåller drive/prefix: {name}"
    normalized = posixpath.normpath(name)
    assert normalized == name, f"ZIP-namn normaliserar bort segment: {name}"
    assert ".." not in Path(name).parts, f"ZIP-namn innehåller parentsegment: {name}"


def assert_data_dir_has_no_plaintext(tokens: list[bytes]) -> None:
    data_root = Path(runtime_dir.name)
    for path in data_root.rglob("*"):
        if not path.is_file():
            continue
        data = path.read_bytes()
        for token in tokens:
            assert token not in data, f"{token!r} finns i klartext på disk i {path}"


def main() -> None:
    try:
        db.init_db()
        with TestClient(app) as client:
            unauthenticated = client.get("/export/zip", follow_redirects=False)
            assert unauthenticated.status_code == 303 and unauthenticated.headers["location"] == "/login"
            api_unauthenticated = client.get("/api/v1/documents")
            assert api_unauthenticated.status_code == 401
            failed_login_page = client.get("/login")
            failed_login = client.post(
                "/login",
                data={"username": "admin", "password": "wrong-password", "csrf_token": csrf(failed_login_page)},
            )
            assert failed_login.status_code == 401
            fail2ban_log = Path(runtime_dir.name) / "logs" / "fail2ban-auth.log"
            assert fail2ban_log.exists() and "LOGIN_FAILED" in fail2ban_log.read_text(encoding="utf-8")

            login_page = client.get("/login")
            assert "Första setup" in login_page.text
            login = client.post(
                "/login",
                data={"username": "admin", "password": "12345", "csrf_token": csrf(login_page)},
                follow_redirects=False,
            )
            assert login.status_code == 303 and login.headers["location"] == "/change-password"
            blocked = client.get("/", follow_redirects=False)
            assert blocked.status_code == 303 and blocked.headers["location"] == "/change-password"

            change_page = client.get("/change-password")
            changed = client.post(
                "/change-password",
                data={
                    "current_password": "12345",
                    "new_password": "acceptance-12345",
                    "confirm_password": "acceptance-12345",
                    "csrf_token": csrf(change_page),
                },
                follow_redirects=False,
            )
            assert changed.status_code == 303 and changed.headers["location"] == "/"
            admin_archive = client.get("/")
            assert admin_archive.status_code == 200
            assert "Adminläget har inte direkt filåtkomst" in admin_archive.text
            login_after_setup = client.get("/login")
            client.post("/logout", data={"csrf_token": csrf(admin_archive)}, follow_redirects=False)
            login_after_setup = client.get("/login")
            assert "Första setup" not in login_after_setup.text
            admin_login_page = client.get("/login")
            admin_login = client.post(
                "/login",
                data={"username": "admin", "password": "acceptance-12345", "csrf_token": csrf(admin_login_page)},
                follow_redirects=False,
            )
            assert admin_login.status_code == 303
            admin_upload = client.get("/upload")
            assert admin_upload.status_code == 403
            admin_token = db.create_api_token(1, "admin-acceptance")
            admin_docs = client.get("/api/v1/documents", headers={"Authorization": f"Bearer {admin_token}"})
            assert admin_docs.status_code == 200 and admin_docs.json() == []
            settings_page = client.get("/settings")
            blocked_invite = client.post(
                "/settings/users/invite",
                data={"csrf_token": csrf(settings_page), "email": "blocked@example.test"},
                follow_redirects=False,
            )
            assert blocked_invite.status_code == 303 and "kr%C3%A4ver" in blocked_invite.headers["location"]
            manual_created = client.post(
                "/settings/users/manual",
                data={
                    "csrf_token": csrf(settings_page),
                    "username": "ArkivUser",
                    "email": "arkivuser@example.test",
                    "temporary_password": "temporary-arkiv-user",
                },
                follow_redirects=False,
            )
            assert manual_created.status_code == 303
            primary_user = db.get_user_by_username("ArkivUser")
            assert primary_user and primary_user["role"] == "user" and primary_user["must_change_password"]
            settings_page = client.get("/settings")
            blocked_reset = client.post(
                f"/settings/users/{primary_user['id']}/reset-email",
                data={"csrf_token": csrf(settings_page)},
                follow_redirects=False,
            )
            assert blocked_reset.status_code == 303 and "SMTP" in blocked_reset.headers["location"]
            blocked_impersonation = client.post(
                f"/settings/impersonate/{primary_user['id']}",
                data={"csrf_token": csrf(settings_page)},
                headers={"x-forwarded-for": "127.0.0.1"},
                follow_redirects=False,
            )
            assert blocked_impersonation.status_code == 403
            allow_saved = client.post(
                "/settings/impersonation",
                data={"csrf_token": csrf(settings_page), "admin_impersonation_allowed_ips": "127.0.0.1"},
                headers={"x-forwarded-for": "127.0.0.1"},
                follow_redirects=False,
            )
            assert allow_saved.status_code == 303
            settings_page = client.get("/settings", headers={"x-forwarded-for": "127.0.0.1"})
            impersonated = client.post(
                f"/settings/impersonate/{primary_user['id']}",
                data={"csrf_token": csrf(settings_page)},
                headers={"x-forwarded-for": "127.0.0.1"},
                follow_redirects=False,
            )
            assert impersonated.status_code == 303
            assert "Impersonerar ArkivUser" in client.get("/").text
            stopped = client.post("/impersonation/stop", data={"csrf_token": csrf(client.get('/'))}, follow_redirects=False)
            assert stopped.status_code == 303
            logged_out_admin = client.post("/logout", data={"csrf_token": csrf(client.get('/settings'))}, follow_redirects=False)
            assert logged_out_admin.status_code == 303
            reset_cli = subprocess.run(
                [sys.executable, "scripts/reset_admin_password.py", "--password-stdin"],
                input="cli-reset-admin\n",
                text=True,
                cwd=PROJECT_ROOT,
                capture_output=True,
                check=False,
            )
            assert reset_cli.returncode == 0, reset_cli.stderr
            assert db.verify_password("cli-reset-admin", db.get_user_by_username("admin")["password_hash"])
            cli_login_page = client.get("/login")
            cli_login = client.post(
                "/login",
                data={"username": "admin", "password": "cli-reset-admin", "csrf_token": csrf(cli_login_page)},
                follow_redirects=False,
            )
            assert cli_login.status_code == 303 and cli_login.headers["location"] == "/change-password"
            cli_change_page = client.get("/change-password")
            cli_changed = client.post(
                "/change-password",
                data={
                    "current_password": "cli-reset-admin",
                    "new_password": "acceptance-12345",
                    "confirm_password": "acceptance-12345",
                    "csrf_token": csrf(cli_change_page),
                },
                follow_redirects=False,
            )
            assert cli_changed.status_code == 303
            client.post("/logout", data={"csrf_token": csrf(client.get('/'))}, follow_redirects=False)
            user_login_page = client.get("/login")
            user_login = client.post(
                "/login",
                data={"username": "ArkivUser", "password": "temporary-arkiv-user", "csrf_token": csrf(user_login_page)},
                follow_redirects=False,
            )
            assert user_login.status_code == 303 and user_login.headers["location"] == "/change-password"
            user_change_page = client.get("/change-password")
            user_changed = client.post(
                "/change-password",
                data={
                    "current_password": "temporary-arkiv-user",
                    "new_password": "acceptance-user",
                    "confirm_password": "acceptance-user",
                    "csrf_token": csrf(user_change_page),
                },
                follow_redirects=False,
            )
            assert user_changed.status_code == 303 and user_changed.headers["location"] == "/"
            primary_user = db.get_user_by_username("ArkivUser")
            assert primary_user and primary_user["role"] == "user" and not primary_user["must_change_password"]
            assert db.verify_password("acceptance-user", primary_user["password_hash"])
            assert len(DOCUMENT_TEMPLATES) >= 77
            assert any(item["id"] == "health_insurance" and item["name"] == "Sjukförsäkring" for item in DOCUMENT_TEMPLATES)
            assert any(item["id"] == "dog_insurance" and item["name"] == "Hundförsäkring" for item in DOCUMENT_TEMPLATES)
            assert any(item["id"] == "securities_account" and item["name"] == "Värdepapperskonto" for item in DOCUMENT_TEMPLATES)
            assert any(item["id"] == "password_vault_export" and item["name"] == "Lösenord/valvexport" for item in DOCUMENT_TEMPLATES)
            assert any(item["id"] == "backup_recovery_plan" and item["name"] == "Backup/återställning" for item in DOCUMENT_TEMPLATES)
            upload_page = client.get("/upload")
            assert upload_page.status_code == 200
            assert 'name="files"' in upload_page.text and 'name="template_id"' in upload_page.text

            token = db.create_api_token(primary_user["id"], "acceptance")
            uploaded_ids: list[int] = []
            multi_upload = client.post(
                "/upload",
                files=[
                    ("files", ("multi-kvitto.txt", b"unikmulti kvitto", "text/plain")),
                    ("files", ("multi-sjuk.txt", b"unikmulti sjukforsakring", "text/plain")),
                ],
                data={"csrf_token": csrf(upload_page), "template_id": ["receipt", "health_insurance"], "tags": "multi,acceptance"},
                follow_redirects=False,
            )
            assert multi_upload.status_code == 303 and multi_upload.headers["location"].startswith("/?message=2")
            archive_after_multi_upload = client.get(multi_upload.headers["location"])
            assert archive_after_multi_upload.status_code == 200
            assert '<div class="doc-tags"' in archive_after_multi_upload.text
            assert "<span>multi</span>" in archive_after_multi_upload.text
            assert "<span>acceptance</span>" in archive_after_multi_upload.text
            multi_docs = client.get(
                "/api/v1/documents",
                headers={"Authorization": f"Bearer {token}"},
                params={"tag": "multi"},
            )
            assert multi_docs.status_code == 200
            multi_by_name = {row["original_filename"]: row for row in multi_docs.json()}
            assert multi_by_name["multi-kvitto.txt"]["template_id"] == "receipt"
            assert multi_by_name["multi-sjuk.txt"]["template_id"] == "health_insurance"
            for filename, content, token_word, must_index in SAMPLES:
                response = client.post(
                    "/api/v1/documents",
                    headers={"Authorization": f"Bearer {token}"},
                    files={"file": (filename, content)},
                    data={"template_id": "receipt", "tags": "acceptance,import"},
                )
                assert response.status_code == 200, f"{filename}: {response.text}"
                payload = response.json()
                uploaded_ids.append(payload["id"])
                assert Path(payload["original_filename"]).name == payload["original_filename"]
                assert (Path(runtime_dir.name) / "uploads").exists()
                if must_index:
                    assert payload["extraction_status"] == "indexed", f"{filename} indexerades inte."
                    search = client.get(
                        "/api/v1/documents",
                        headers={"Authorization": f"Bearer {token}"},
                        params={"q": token_word},
                    )
                    assert search.status_code == 200
                    assert any(row["id"] == payload["id"] for row in search.json()), f"{filename} saknas i sök."
                else:
                    assert payload["extraction_status"] in {"indexed", "archived_only"}
                    assert "width" in payload["metadata"] and "height" in payload["metadata"]

            traversal = client.post(
                "/api/v1/documents",
                headers={"Authorization": f"Bearer {token}"},
                files={"file": ("../../hemligt\\konto.txt", b"uniktraversal 070-123 45 67 konto\n")},
            )
            assert traversal.status_code == 200
            traversal_payload = traversal.json()
            assert traversal_payload["original_filename"] == "konto.txt"
            delete_page = client.get(f"/documents/{traversal_payload['id']}")
            assert delete_page.status_code == 200
            assert f'action="/documents/{traversal_payload["id"]}/delete"' in delete_page.text
            with db.connect() as conn:
                delete_row = conn.execute(
                    "SELECT storage_path, text_path FROM documents WHERE id = ?",
                    (traversal_payload["id"],),
                ).fetchone()
            delete_paths = [Path(delete_row["storage_path"]), Path(delete_row["text_path"])]
            delete_invalid_csrf = client.post(
                f"/documents/{traversal_payload['id']}/delete",
                data={"csrf_token": "fel-token"},
                follow_redirects=False,
            )
            assert delete_invalid_csrf.status_code == 403
            deleted = client.post(
                f"/documents/{traversal_payload['id']}/delete",
                data={"csrf_token": csrf(delete_page)},
                follow_redirects=False,
            )
            assert deleted.status_code == 303 and deleted.headers["location"].startswith("/?message=")
            assert client.get(f"/documents/{traversal_payload['id']}").status_code == 404
            after_delete_docs = client.get("/api/v1/documents", headers={"Authorization": f"Bearer {token}"})
            assert all(row["id"] != traversal_payload["id"] for row in after_delete_docs.json())
            assert all(not path.exists() for path in delete_paths)

            detail_page = client.get(f"/documents/{uploaded_ids[0]}")
            assert detail_page.status_code == 200
            assert f'action="/documents/{uploaded_ids[0]}/classification"' in detail_page.text
            assert "Checksummor verifierade." in detail_page.text
            assert "Krypterad MD5" in detail_page.text
            invalid_csrf = client.post(
                f"/documents/{uploaded_ids[0]}/classification",
                data={"csrf_token": "fel-token", "template_id": "health_insurance", "tags": "redigerad,acceptance"},
                follow_redirects=False,
            )
            assert invalid_csrf.status_code == 403
            updated_classification = client.post(
                f"/documents/{uploaded_ids[0]}/classification",
                data={"csrf_token": csrf(detail_page), "template_id": "health_insurance", "tags": "redigerad,acceptance"},
                follow_redirects=False,
            )
            assert updated_classification.status_code == 303
            assert updated_classification.headers["location"] == f"/documents/{uploaded_ids[0]}?message=classification-updated"
            updated_detail = client.get(updated_classification.headers["location"])
            assert updated_detail.status_code == 200
            assert "Mall och taggar uppdaterade." in updated_detail.text
            assert "Sjukförsäkring" in updated_detail.text
            assert 'value="redigerad,acceptance"' in updated_detail.text
            old_tag_search = client.get(
                "/api/v1/documents",
                headers={"Authorization": f"Bearer {token}"},
                params={"tag": "import"},
            )
            assert old_tag_search.status_code == 200
            assert all(row["id"] != uploaded_ids[0] for row in old_tag_search.json())
            new_tag_search = client.get(
                "/api/v1/documents",
                headers={"Authorization": f"Bearer {token}"},
                params={"tag": "redigerad"},
            )
            assert new_tag_search.status_code == 200
            assert any(row["id"] == uploaded_ids[0] for row in new_tag_search.json())
            api_documents = client.get("/api/v1/documents", headers={"Authorization": f"Bearer {token}"})
            assert api_documents.status_code == 200
            edited_api_doc = next(row for row in api_documents.json() if row["id"] == uploaded_ids[0])
            assert edited_api_doc["tags"] == "redigerad,acceptance"
            assert edited_api_doc["template_id"] == "health_insurance"
            assert len(edited_api_doc["md5_plain"]) == 32
            assert len(edited_api_doc["sha256_encrypted"]) == 64
            assert len(edited_api_doc["md5_encrypted"]) == 32
            metadata_after_tag_edit = client.get("/export/metadata.json")
            assert metadata_after_tag_edit.status_code == 200
            edited_export_doc = next(row for row in metadata_after_tag_edit.json() if row["id"] == uploaded_ids[0])
            assert edited_export_doc["tags"] == "redigerad,acceptance"
            assert edited_export_doc["template_id"] == "health_insurance"

            forbidden = client.post(
                "/api/v1/documents",
                headers={"Authorization": f"Bearer {token}"},
                files={"file": ("malware.exe", b"nej")},
            )
            assert forbidden.status_code == 400

            too_large = client.post(
                "/api/v1/documents",
                headers={"Authorization": f"Bearer {token}"},
                files={"file": ("large.txt", b"x" * 1_000_001)},
            )
            assert too_large.status_code == 413

            chat_page = client.get("/chat")
            chat = client.post(
                "/chat",
                data={
                    "csrf_token": csrf(chat_page),
                    "question": "Vad står i dokumentet?",
                    "document_ids": str(uploaded_ids[0]),
                    "redact": "on",
                },
            )
            assert chat.status_code == 400
            assert "AI är inte aktiverat" in chat.text
            assert "19800101-1234" not in chat.text
            assert "[MASKAT_PERSONNUMMER]" in chat.text
            assert redact_text("me@example.test +46 70 123 45 67 SE4550000000058398257466").count("[MASK") >= 3

            metadata_json = client.get("/export/metadata.json")
            assert metadata_json.status_code == 200
            assert len(metadata_json.json()) >= len(uploaded_ids)
            metadata_csv = client.get("/export/metadata.csv")
            assert metadata_csv.status_code == 200 and "original_filename" in metadata_csv.text

            exported = client.get("/export/zip", params={"ids": ",".join(map(str, uploaded_ids))})
            assert exported.status_code == 200
            with zipfile.ZipFile(io.BytesIO(exported.content)) as archive:
                names = archive.namelist()
                assert "manifest.json" in names
                assert any(name.startswith("original/") and name.endswith("text-token.txt") for name in names)
                for name in names:
                    assert_zip_member_safe(name)
                manifest = json.loads(archive.read("manifest.json"))
                assert any(item["original_filename"] == "text-token.txt" for item in manifest)

            assert safe_zip_member("original", "../../..\\CON.txt", 99) == "original/99_CON.txt"

            bob_invite = db.create_user_invite("bob@example.test", primary_user["id"])
            bob_id = db.accept_user_invite(bob_invite, "Bob", "acceptance-bob")
            assert bob_id
            bob_token = db.create_api_token(bob_id, "bob")
            bob_before_share = client.get("/api/v1/documents", headers={"Authorization": f"Bearer {bob_token}"})
            assert bob_before_share.status_code == 200
            assert all(row["id"] != uploaded_ids[0] for row in bob_before_share.json())
            sent_mail: list[str] = []
            old_env = {key: os.environ.get(key) for key in ["SMTP_HOST", "SMTP_PORT", "SMTP_USER", "SMTP_PASS", "MAIL_FROM"]}
            os.environ.update(
                {
                    "SMTP_HOST": "smtp.example.test",
                    "SMTP_PORT": "587",
                    "SMTP_USER": "noreply@example.test",
                    "SMTP_PASS": "secret",
                    "MAIL_FROM": "noreply@example.test",
                }
            )
            from app.services import mail as app_mail

            real_send_mail = app_mail.send_mail
            app_mail.send_mail = lambda to_addr, subject, text, **kwargs: sent_mail.append(text)  # type: ignore[assignment]
            try:
                share_page = client.get(f"/documents/{uploaded_ids[0]}")
                share_response = client.post(
                    f"/documents/{uploaded_ids[0]}/share",
                    data={"csrf_token": csrf(share_page), "recipient_email": "bob@example.test"},
                    follow_redirects=False,
                )
                assert share_response.status_code == 303
                assert sent_mail and "/invites/shares/" in sent_mail[-1]
                share_token = sent_mail[-1].rsplit("/invites/shares/", 1)[1].strip().split()[0]
            finally:
                app_mail.send_mail = real_send_mail  # type: ignore[assignment]
                for key, value in old_env.items():
                    if value is None:
                        os.environ.pop(key, None)
                    else:
                        os.environ[key] = value
            client.post("/logout", data={"csrf_token": csrf(share_page)}, follow_redirects=False)
            bob_login_page = client.get("/login")
            bob_login = client.post(
                "/login",
                data={"username": "Bob", "password": "acceptance-bob", "csrf_token": csrf(bob_login_page)},
                follow_redirects=False,
            )
            assert bob_login.status_code == 303
            accepted_share = client.get(f"/invites/shares/{share_token}", follow_redirects=False)
            assert accepted_share.status_code == 303
            bob_after_share = client.get("/api/v1/documents", headers={"Authorization": f"Bearer {bob_token}"})
            assert any(row["id"] == uploaded_ids[0] for row in bob_after_share.json())
            bob_delete_page = client.get(f"/documents/{uploaded_ids[0]}")
            assert f'action="/documents/{uploaded_ids[0]}/delete"' not in bob_delete_page.text
            bob_download = client.get(f"/documents/{uploaded_ids[0]}/download")
            assert bob_download.status_code == 200 and b"uniktext alfa" in bob_download.content
            bob_export = client.get("/export/zip", params={"ids": str(uploaded_ids[0])})
            assert bob_export.status_code == 200
            bob_delete = client.post(
                f"/documents/{uploaded_ids[0]}/delete",
                data={"csrf_token": csrf(bob_delete_page)},
                follow_redirects=False,
            )
            assert bob_delete.status_code == 404
            bob_update = client.post(
                f"/documents/{uploaded_ids[0]}/classification",
                data={"csrf_token": csrf(bob_delete_page), "template_id": "receipt", "tags": "nope"},
                follow_redirects=False,
            )
            assert bob_update.status_code == 404
            client.post("/logout", data={"csrf_token": csrf(bob_delete_page)}, follow_redirects=False)
            user_login_page = client.get("/login")
            user_login = client.post(
                "/login",
                data={"username": "ArkivUser", "password": "acceptance-user", "csrf_token": csrf(user_login_page)},
                follow_redirects=False,
            )
            assert user_login.status_code == 303
            revoke_page = client.get(f"/documents/{uploaded_ids[0]}")
            revoked = client.post(
                f"/documents/{uploaded_ids[0]}/share/{bob_id}/revoke",
                data={"csrf_token": csrf(revoke_page)},
                follow_redirects=False,
            )
            assert revoked.status_code == 303
            bob_after_revoke = client.get("/api/v1/documents", headers={"Authorization": f"Bearer {bob_token}"})
            assert all(row["id"] != uploaded_ids[0] for row in bob_after_revoke.json())

            client.post("/logout", data={"csrf_token": csrf(revoke_page)}, follow_redirects=False)
            admin_login_page = client.get("/login")
            admin_login = client.post(
                "/login",
                data={"username": "admin", "password": "acceptance-12345", "csrf_token": csrf(admin_login_page)},
                follow_redirects=False,
            )
            assert admin_login.status_code == 303

            settings_page = client.get("/settings")
            assert settings_page.status_code == 200
            assert "Disabled" in settings_page.text
            token_form = client.post(
                "/settings/api-token",
                data={"csrf_token": csrf(settings_page), "token_name": "acceptance-ui"},
                follow_redirects=False,
            )
            assert token_form.status_code == 303
            token_location = token_form.headers["location"]
            assert token_location.startswith("/settings?token_created=")
            assert "new_token=" not in token_location and "dk_" not in token_location
            token_page = client.get(token_location)
            assert token_page.status_code == 200
            token_match = re.search(r"<code>(dk_[^<]+)</code>", token_page.text)
            assert token_match, "Ny API-token ska visas en gång efter skapande."
            ui_token = html.unescape(token_match.group(1))
            assert db.authenticate_token(ui_token), "Skapad UI-token ska fungera som bearer token."
            replayed_token_page = client.get(token_location)
            assert ui_token not in replayed_token_page.text, "API-token ska inte visas igen via samma flash-länk."

            reset_mail: list[str] = []
            reset_old_env = {key: os.environ.get(key) for key in ["SMTP_HOST", "SMTP_PORT", "SMTP_USER", "SMTP_PASS", "MAIL_FROM"]}
            os.environ.update(
                {
                    "SMTP_HOST": "smtp.example.test",
                    "SMTP_PORT": "587",
                    "SMTP_USER": "noreply@example.test",
                    "SMTP_PASS": "secret",
                    "MAIL_FROM": "noreply@example.test",
                }
            )
            from app.services import mail as app_mail

            real_send_mail = app_mail.send_mail
            app_mail.send_mail = lambda to_addr, subject, text, **kwargs: reset_mail.append(text)  # type: ignore[assignment]
            try:
                reset_settings_page = client.get("/settings")
                reset_response = client.post(
                    f"/settings/users/{primary_user['id']}/reset-email",
                    data={"csrf_token": csrf(reset_settings_page)},
                    follow_redirects=False,
                )
                assert reset_response.status_code == 303
                assert reset_mail and "/invites/password/" in reset_mail[-1]
                reset_token = reset_mail[-1].rsplit("/invites/password/", 1)[1].strip().split()[0]
            finally:
                app_mail.send_mail = real_send_mail  # type: ignore[assignment]
                for key, value in reset_old_env.items():
                    if value is None:
                        os.environ.pop(key, None)
                    else:
                        os.environ[key] = value
            reset_page = client.get(f"/invites/password/{reset_token}")
            reset_done = client.post(
                f"/invites/password/{reset_token}",
                data={
                    "csrf_token": csrf(reset_page),
                    "password": "acceptance-user-reset",
                    "confirm_password": "acceptance-user-reset",
                },
                follow_redirects=False,
            )
            assert reset_done.status_code == 303
            assert db.verify_password("acceptance-user-reset", db.get_user(primary_user["id"])["password_hash"])
            client.post("/logout", data={"csrf_token": csrf(client.get('/'))}, follow_redirects=False)
            admin_login_page = client.get("/login")
            admin_login = client.post(
                "/login",
                data={"username": "admin", "password": "acceptance-12345", "csrf_token": csrf(admin_login_page)},
                follow_redirects=False,
            )
            assert admin_login.status_code == 303

            provider_forms = [
                {
                    "ai_provider": "openai",
                    "ai_openai_api_key": "test-openai-key",
                    "ai_openai_model": "gpt-test",
                    "ai_openai_base_url": "http://127.0.0.1:9/v1",
                },
                {
                    "ai_provider": "claude",
                    "ai_claude_api_key": "test-claude-key",
                    "ai_claude_model": "claude-test",
                },
                {
                    "ai_provider": "ollama",
                    "ai_ollama_base_url": "http://127.0.0.1:11434",
                    "ai_ollama_model": "llama-test",
                },
            ]
            for form in provider_forms:
                settings_page = client.get("/settings")
                ai_saved = client.post(
                    "/settings/ai",
                    data={"csrf_token": csrf(settings_page), **form},
                    follow_redirects=False,
                )
                assert ai_saved.status_code == 303
                settings = db.get_settings()
                assert settings["ai_provider"] == form["ai_provider"]
                assert settings["ai_enabled"] == "false" and settings["ai_last_test_ok"] == "false"
                masked_page = client.get("/settings")
                assert "test-openai-key" not in masked_page.text
                assert "test-claude-key" not in masked_page.text

            mail_settings_page = client.get("/settings")
            mail_saved = client.post(
                "/settings/mail-import",
                data={
                    "csrf_token": csrf(mail_settings_page),
                    "mail_import_enabled": "on",
                    "mail_import_protocol": "imap",
                    "mail_import_host": "mail.example.test",
                    "mail_import_port": "993",
                    "mail_import_ssl": "on",
                    "mail_import_username": "archive@example.test",
                    "mail_import_password": "test-mail-password",
                    "mail_import_folder": "INBOX",
                    "mail_import_delete_after_handled": "on",
                    "mail_import_poll_interval_seconds": "120",
                    "mail_import_max_messages": "5",
                    "mail_import_min_inline_image_bytes": "10240",
                    "mail_import_import_eml_without_attachments": "on",
                    "mail_import_default_tags": "mailimport,acceptance",
                },
                follow_redirects=False,
            )
            assert mail_saved.status_code == 303
            mail_settings = db.get_settings()
            assert mail_settings["mail_import_protocol"] == "imap"
            assert mail_settings["mail_import_delete_after_handled"] == "true"
            assert mail_settings["mail_import_password"] == "test-mail-password"
            masked_mail_settings = client.get("/settings")
            assert "test-mail-password" not in masked_mail_settings.text

            attachment_mail = EmailMessage()
            attachment_mail["From"] = "archive@forsakring.example"
            attachment_mail["To"] = "dokumenteraren@example.test"
            attachment_mail["Subject"] = "Hundförsäkring veterinär kvitto"
            attachment_mail.set_content("Se bifogad hundförsäkring.")
            attachment_mail.add_attachment(
                b"unikmail hund forsakring veterinarkvitto",
                maintype="text",
                subtype="plain",
                filename="hundforsakring-kvitto.txt",
            )
            mail_results = mail_importer.process_message_bytes(attachment_mail.as_bytes(), mail_settings, uid="acceptance-uid-1")
            assert any(item["status"] == "imported" for item in mail_results)
            mail_search = client.get(
                "/api/v1/documents",
                headers={"Authorization": f"Bearer {token}"},
                params={"q": "unikmail"},
            )
            assert mail_search.status_code == 200
            mail_doc = next(row for row in mail_search.json() if row["original_filename"] == "hundforsakring-kvitto.txt")
            assert mail_doc["template_id"] == "dog_insurance"
            assert "automatiskt sorterad" in mail_doc["tags"]
            assert "mailimport" in mail_doc["tags"]
            assert "from:forsakring.example" in mail_doc["tags"]

            no_attachment_mail = EmailMessage()
            no_attachment_mail["From"] = "ops@example.test"
            no_attachment_mail["To"] = "dokumenteraren@example.test"
            no_attachment_mail["Subject"] = "Recovery codes"
            no_attachment_mail.set_content("unikemlimport backup codes")
            eml_results = mail_importer.process_message_bytes(no_attachment_mail.as_bytes(), mail_settings, uid="acceptance-uid-2")
            assert any(item["status"] == "imported" for item in eml_results)
            eml_search = client.get(
                "/api/v1/documents",
                headers={"Authorization": f"Bearer {token}"},
                params={"q": "unikemlimport"},
            )
            assert eml_search.status_code == 200
            assert any(row["original_filename"].endswith(".eml") for row in eml_search.json())

            import_dir = Path(os.environ["IMPORT_DIR"])
            import_dir.mkdir(exist_ok=True)
            db.init_db()
            assert not import_dir.is_relative_to(Path(runtime_dir.name)), "Importkatalogen ska ligga utanför DATA_DIR."
            assert import_dir.stat().st_mode & 0o002, "Importkatalogen ska vara skrivbar från hosten."
            imported_file = import_dir / "slukad-import.txt"
            imported_file.write_bytes(b"unikimport hemlig landing zone")
            importer.process_import_once()
            import_results = importer.process_import_once()
            assert any(item["status"] == "imported" for item in import_results)
            assert not imported_file.exists(), "Lyckad import ska tas bort från importkatalogen."
            imported_search = client.get(
                "/api/v1/documents",
                headers={"Authorization": f"Bearer {token}"},
                params={"q": "unikimport"},
            )
            assert imported_search.status_code == 200
            assert any(row["original_filename"] == "slukad-import.txt" for row in imported_search.json())

            blocked_import = import_dir / "blockerad.exe"
            blocked_import.write_bytes(b"unikblockerad plaintext ska inte ligga kvar")
            importer.process_import_once()
            failed_results = importer.process_import_once()
            assert any(item["status"] == "failed" for item in failed_results)
            assert not blocked_import.exists(), "Misslyckad import ska inte ligga kvar i importkatalogen."
            failed_files = list((Path(runtime_dir.name) / "import_failed").glob("*.enc"))
            assert failed_files and failed_files[-1].read_bytes().startswith(ENC_PREFIX.encode("ascii"))

            export_artifact = Path(runtime_dir.name) / "exports" / "dokumenteraren_export.zip.enc"
            assert export_artifact.exists() and export_artifact.read_bytes().startswith(ENC_PREFIX.encode("ascii"))

            assert_data_dir_has_no_plaintext(
                [
                    b"uniktext alfa",
                    b"unikdocx avtal",
                    b"unikxlsx kvitto",
                    b"unikimport hemlig",
                    b"unikblockerad plaintext",
                    "redigerad,acceptance".encode("utf-8"),
                ]
            )

        print(f"ACCEPTANCE_OK data_dir={runtime_dir.name}")
    finally:
        runtime_dir.cleanup()


if __name__ == "__main__":
    main()
