from __future__ import annotations

import csv
import email
import json
import mimetypes
import shutil
import subprocess
import tempfile
from dataclasses import dataclass
from email import policy
from html.parser import HTMLParser
from pathlib import Path
from typing import Any

import yaml
from docx import Document as DocxDocument
from odf import teletype
from odf.opendocument import load as odf_load
from odf.text import P as OdfParagraph
from openpyxl import load_workbook
from PIL import Image
from pdf2image import convert_from_path
from pypdf import PdfReader
from pptx import Presentation
from striprtf.striprtf import rtf_to_text


@dataclass
class ExtractionResult:
    text: str
    metadata: dict[str, Any]
    status: str


class TextHTMLParser(HTMLParser):
    def __init__(self) -> None:
        super().__init__()
        self.parts: list[str] = []

    def handle_data(self, data: str) -> None:
        value = data.strip()
        if value:
            self.parts.append(value)

    def text(self) -> str:
        return "\n".join(self.parts)


def extract_document(path: Path, extension: str, original_filename: str) -> ExtractionResult:
    metadata: dict[str, Any] = {
        "original_filename": original_filename,
        "extension": extension,
        "mime_guess": mimetypes.guess_type(original_filename)[0] or "application/octet-stream",
    }
    try:
        if extension in {"txt", "md", "ini", "conf", "log"}:
            return ExtractionResult(read_text(path), metadata, "indexed")
        if extension == "rtf":
            return ExtractionResult(rtf_to_text(read_text(path)), metadata, "indexed")
        if extension in {"csv", "tsv"}:
            return extract_csv(path, metadata, "\t" if extension == "tsv" else ",")
        if extension in {"json", "yaml", "yml", "xml"}:
            return extract_structured_text(path, metadata, extension)
        if extension in {"html", "htm"}:
            parser = TextHTMLParser()
            parser.feed(read_text(path))
            return ExtractionResult(parser.text(), metadata, "indexed")
        if extension == "eml":
            return extract_eml(path, metadata)
        if extension == "pdf":
            return extract_pdf(path, metadata)
        if extension == "docx":
            return extract_docx(path, metadata)
        if extension == "xlsx":
            return extract_xlsx(path, metadata)
        if extension == "pptx":
            return extract_pptx(path, metadata)
        if extension in {"odt", "ods", "odp", "ott", "ots", "otp"}:
            return extract_odf(path, metadata)
        if extension in {"jpg", "jpeg", "png", "webp", "tif", "tiff", "bmp", "gif", "heic"}:
            return extract_image(path, metadata)
        if extension in {"doc", "xls", "ppt"}:
            return extract_legacy_office(path, metadata)
    except Exception as exc:
        metadata["extraction_error"] = exc.__class__.__name__
        return ExtractionResult("", metadata, "archived_only")
    return ExtractionResult("", metadata, "archived_only")


def read_text(path: Path) -> str:
    data = path.read_bytes()
    for encoding in ("utf-8", "utf-8-sig", "latin-1"):
        try:
            return data.decode(encoding)
        except UnicodeDecodeError:
            continue
    return data.decode("utf-8", errors="replace")


def compact(text: str, limit: int = 250_000) -> str:
    text = "\n".join(line.rstrip() for line in text.splitlines())
    text = "\n".join(line for line in text.splitlines() if line.strip())
    return text[:limit]


def extract_csv(path: Path, metadata: dict[str, Any], delimiter: str) -> ExtractionResult:
    rows: list[str] = []
    with path.open("r", encoding="utf-8", errors="replace", newline="") as handle:
        reader = csv.reader(handle, delimiter=delimiter)
        for index, row in enumerate(reader):
            if index < 200:
                rows.append(" | ".join(row))
            if index == 0:
                metadata["columns"] = row
            metadata["row_count_seen"] = index + 1
    return ExtractionResult(compact("\n".join(rows)), metadata, "indexed")


def extract_structured_text(path: Path, metadata: dict[str, Any], extension: str) -> ExtractionResult:
    text = read_text(path)
    if extension == "json":
        parsed = json.loads(text)
        metadata["top_level_type"] = type(parsed).__name__
    elif extension in {"yaml", "yml"}:
        parsed = yaml.safe_load(text)
        metadata["top_level_type"] = type(parsed).__name__
    return ExtractionResult(compact(text), metadata, "indexed")


def extract_eml(path: Path, metadata: dict[str, Any]) -> ExtractionResult:
    msg = email.message_from_bytes(path.read_bytes(), policy=policy.default)
    metadata.update(
        {
            "subject": msg.get("subject", ""),
            "from": msg.get("from", ""),
            "to": msg.get("to", ""),
            "date": msg.get("date", ""),
        }
    )
    parts: list[str] = []
    if msg.is_multipart():
        for part in msg.walk():
            if part.get_content_type() == "text/plain":
                parts.append(part.get_content())
    elif msg.get_content_type() == "text/plain":
        parts.append(msg.get_content())
    return ExtractionResult(compact("\n".join(parts)), metadata, "indexed")


def extract_pdf(path: Path, metadata: dict[str, Any]) -> ExtractionResult:
    reader = PdfReader(str(path))
    metadata["pages"] = len(reader.pages)
    if reader.metadata:
        metadata["pdf_metadata"] = {str(k): str(v) for k, v in reader.metadata.items()}
    text = "\n".join(page.extract_text() or "" for page in reader.pages)
    if text.strip():
        return ExtractionResult(compact(text), metadata, "indexed")
    ocr_text = ocr_pdf(path, metadata)
    if ocr_text.strip():
        return ExtractionResult(compact(ocr_text), metadata, "indexed")
    return ExtractionResult("", metadata, "archived_only")


def ocr_pdf(path: Path, metadata: dict[str, Any]) -> str:
    if not shutil.which("tesseract") or not shutil.which("pdftoppm"):
        metadata["pdf_ocr"] = "missing_tesseract_or_poppler"
        return ""
    try:
        images = convert_from_path(str(path), dpi=200, first_page=1, last_page=20)
    except Exception as exc:
        metadata["pdf_ocr_error"] = exc.__class__.__name__
        return ""
    parts = [ocr_image(image) for image in images]
    metadata["pdf_ocr_pages"] = len(images)
    return "\n".join(parts)


def extract_docx(path: Path, metadata: dict[str, Any]) -> ExtractionResult:
    doc = DocxDocument(str(path))
    props = doc.core_properties
    metadata["author"] = props.author or ""
    metadata["created"] = props.created.isoformat() if props.created else ""
    text = "\n".join(p.text for p in doc.paragraphs if p.text.strip())
    return ExtractionResult(compact(text), metadata, "indexed" if text.strip() else "archived_only")


def extract_xlsx(path: Path, metadata: dict[str, Any]) -> ExtractionResult:
    workbook = load_workbook(str(path), read_only=True, data_only=True)
    metadata["sheets"] = workbook.sheetnames
    parts: list[str] = []
    for sheet in workbook.worksheets:
        parts.append(f"# {sheet.title}")
        for idx, row in enumerate(sheet.iter_rows(values_only=True)):
            if idx >= 200:
                break
            values = ["" if value is None else str(value) for value in row]
            if any(values):
                parts.append(" | ".join(values))
    text = "\n".join(parts)
    return ExtractionResult(compact(text), metadata, "indexed" if text.strip() else "archived_only")


def extract_pptx(path: Path, metadata: dict[str, Any]) -> ExtractionResult:
    presentation = Presentation(str(path))
    metadata["slides"] = len(presentation.slides)
    parts: list[str] = []
    for slide in presentation.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                parts.append(shape.text)
    text = "\n".join(parts)
    return ExtractionResult(compact(text), metadata, "indexed" if text.strip() else "archived_only")


def extract_odf(path: Path, metadata: dict[str, Any]) -> ExtractionResult:
    doc = odf_load(str(path))
    metadata["odf_mimetype"] = getattr(doc, "mimetype", "")
    text = teletype.extractText(doc.text) if getattr(doc, "text", None) else ""
    if not text.strip():
        text = "\n".join(
            part
            for part in (teletype.extractText(paragraph) for paragraph in doc.getElementsByType(OdfParagraph))
            if part.strip()
        )
    return ExtractionResult(compact(text), metadata, "indexed" if text.strip() else "archived_only")


def extract_image(path: Path, metadata: dict[str, Any]) -> ExtractionResult:
    with Image.open(path) as image:
        metadata["image_format"] = image.format
        metadata["width"] = image.width
        metadata["height"] = image.height
        metadata["mode"] = image.mode
        exif = image.getexif()
        if exif:
            metadata["exif_keys"] = [str(key) for key in list(exif.keys())[:30]]
        text = ocr_image(image)
    return ExtractionResult(compact(text), metadata, "indexed" if text.strip() else "archived_only")


def ocr_image(image: Image.Image) -> str:
    if not shutil.which("tesseract"):
        return ""
    try:
        import pytesseract

        return pytesseract.image_to_string(image, lang="swe+eng")
    except Exception:
        return ""


def extract_legacy_office(path: Path, metadata: dict[str, Any]) -> ExtractionResult:
    if not shutil.which("soffice"):
        metadata["legacy_office_converter"] = "missing"
        return ExtractionResult("", metadata, "archived_only")
    with tempfile.TemporaryDirectory() as tmp:
        outdir = Path(tmp)
        result = subprocess.run(
            ["soffice", "--headless", "--convert-to", "txt:Text", "--outdir", str(outdir), str(path)],
            check=False,
            stdout=subprocess.DEVNULL,
            stderr=subprocess.DEVNULL,
            timeout=30,
        )
        metadata["legacy_office_converter_exit"] = result.returncode
        converted = list(outdir.glob("*.txt"))
        if converted:
            text = read_text(converted[0])
            return ExtractionResult(compact(text), metadata, "indexed" if text.strip() else "archived_only")
    return ExtractionResult("", metadata, "archived_only")
